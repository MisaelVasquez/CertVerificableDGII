#!/usr/bin/env python3
"""Genera la presentación ejecutiva del caso de uso (DGII · credencial verificable)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Paleta institucional DGII, atenuada para reducir brillo/deslumbramiento en pantalla.
GREEN=RGBColor(0x2c,0x6d,0x1d); GREEN_D=RGBColor(0x1c,0x4d,0x12); LIME=RGBColor(0x74,0xa5,0x2e)
NAVY=RGBColor(0x12,0x30,0x56); RED=RGBColor(0xb0,0x28,0x1d); GOLD=RGBColor(0xa6,0x7c,0x2b)
INK=RGBColor(0x20,0x2a,0x24); GRAY=RGBColor(0x5b,0x66,0x5e); LIGHT=RGBColor(0xee,0xf2,0xe8)
LIGHTN=RGBColor(0xea,0xee,0xf5); WHITE=RGBColor(0xff,0xff,0xff)
F="Calibri"; FL="Calibri Light"

LOGO=os.path.join(os.path.dirname(os.path.abspath(__file__)),"ofv-api","static","logo.png")
LOGO_AR=900/227  # relación de aspecto del logo DGII (ancho/alto)

def logo(s,x,y,w,chip=None):
    """Coloca el logo DGII con ancho w; si chip=color, dibuja una tarjeta de fondo (para fondos oscuros)."""
    h=w/LOGO_AR
    if chip is not None:
        pad=0.16
        rect(s,x-pad,y-pad,w+2*pad,h+2*pad,chip,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    s.shapes.add_picture(LOGO,Inches(x),Inches(y),width=Inches(w))
    return h

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]; SW=13.333

def rect(s,x,y,w,h,color,shape=MSO_SHAPE.RECTANGLE,line=None):
    sp=s.shapes.add_shape(shape,Inches(x),Inches(y),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1)
    sp.shadow.inherit=False
    return sp

def tbox(s,x,y,w,h,anchor=MSO_ANCHOR.TOP):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for m in ("margin_left","margin_right","margin_top","margin_bottom"): setattr(tf,m,0)
    return tf

def run(p,t,size,color,bold=False,font=F,italic=False):
    r=p.add_run(); r.text=t; r.font.size=Pt(size); r.font.color.rgb=color
    r.font.bold=bold; r.font.italic=italic; r.font.name=font; return r

def para(tf,first=False,align=PP_ALIGN.LEFT,space=6,lh=1.12):
    p=tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment=align; p.space_after=Pt(space); p.line_spacing=lh; return p

def header(s,kicker,title,rule=True):
    rect(s,0,0,0.22,7.5,GREEN)                      # lomo verde
    logo(s,10.95,0.46,1.55)                         # logo DGII, esquina superior derecha
    tf=tbox(s,0.85,0.5,9.7,0.4)
    run(para(tf,first=True,space=0),kicker.upper(),12.5,GREEN,bold=True)
    tf=tbox(s,0.85,0.82,9.9,1.0)
    run(para(tf,first=True,space=0,lh=1.02),title,30,NAVY,bold=True,font=FL)
    if rule: rect(s,0.87,1.72,3.2,0.028,LIME)

def footer(s,n):
    tf=tbox(s,0.85,7.02,9,0.3)
    run(para(tf,first=True,space=0),"DGII · Certificación de Impuestos al Día — Credencial Verificable · Bootcamp CDPI",9,GRAY)
    tf=tbox(s,12.0,7.02,0.9,0.3)
    run(para(tf,first=True,space=0,align=PP_ALIGN.RIGHT),str(n),9,GRAY)

def bullets(s,x,y,w,items,size=18,gap=12,color=INK):
    tf=tbox(s,x,y,w,6)
    for i,it in enumerate(items):
        p=para(tf,first=(i==0),space=gap,lh=1.12)
        run(p,"▪  ",size,GREEN,bold=True)
        if isinstance(it,tuple):
            run(p,it[0],size,NAVY,bold=True); run(p,it[1],size,color)
        else:
            run(p,it,size,color)

def card(s,x,y,w,h,accent,title,lines,tsize=17,bsize=13.5):
    rect(s,x,y,w,h,LIGHT,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s,x,y,0.12,h,accent)                        # barra de acento
    tf=tbox(s,x+0.32,y+0.28,w-0.55,0.5)
    run(para(tf,first=True,space=0),title,tsize,accent,bold=True)
    tf=tbox(s,x+0.32,y+0.92,w-0.55,h-1.1)
    for i,ln in enumerate(lines):
        p=para(tf,first=(i==0),space=8,lh=1.12)
        run(p,"·  ",bsize,accent,bold=True); run(p,ln,bsize,INK)

n=[0]
def num(): n[0]+=1; return n[0]

# ---- 1. Portada ----
s=prs.slides.add_slide(BLANK)
rect(s,0,0,SW,7.5,WHITE)
rect(s,0,0,SW,2.35,NAVY)
rect(s,0,2.35,SW,0.12,LIME)
rect(s,0,2.47,SW,0.06,GREEN)
logo(s,9.75,0.62,2.7,chip=WHITE)                # logo DGII sobre chip blanco (banda navy)
tf=tbox(s,0.9,0.7,8.4,0.4); run(para(tf,first=True,space=0),"DIRECCIÓN GENERAL DE IMPUESTOS INTERNOS · BOOTCAMP CDPI",13,LIME,bold=True)
tf=tbox(s,0.9,1.15,11.6,1.1); run(para(tf,first=True,space=0,lh=1.0),"Certificación de Impuestos al Día",34,WHITE,bold=True,font=FL)
tf=tbox(s,0.9,3.1,11.5,1.0); run(para(tf,first=True,space=0,lh=1.05),"Una certificación fiscal como credencial verificable",26,NAVY,bold=True,font=FL)
tf=tbox(s,0.9,4.05,11.0,1.2)
run(para(tf,first=True,space=0,lh=1.2),"Digital, verificable al instante y con privacidad por diseño. El ciudadano prueba que está al día sin exponer más de lo necesario, y la DGII conserva el control mediante revocación en tiempo real.",16,GRAY)
tf=tbox(s,0.9,6.4,11.5,0.4); run(para(tf,first=True,space=0),"Demostración de caso de uso · Prueba de concepto",13,GREEN,bold=True)

# ---- 2. El reto hoy (Contexto y problema · Issuance/Verificación AS-IS) ----
s=prs.slides.add_slide(BLANK); header(s,"Contexto y problema","Cómo se emite y verifica hoy")
tf=tbox(s,0.9,1.92,11.4,0.8)
run(para(tf,first=True,space=0,lh=1.15),"Hoy ya existe un servicio digital eficiente: se solicita en la Oficina Virtual y se emite en minutos (RD$300, PDF firmado con QR validable). El reto no está en emitir, sino en el modelo.",15,GRAY,italic=True)
card(s,0.9,2.86,5.55,1.95,GOLD,"Costo y trámite recurrentes",
     ["Vigencia de 30 días: al caducar, repetir todo y pagar otros RD$300.",
      "Si lo pide un tercero: carta de autorización, cédulas y formulario."])
card(s,6.9,2.86,5.55,1.95,GOLD,"Atado a la DGII en línea",
     ["Validar exige conectividad al portal o app DGII; si falla, se detiene.",
      "Sin escanear el QR, se confía en un PDF fácil de imitar."])
card(s,0.9,4.92,5.55,1.95,GOLD,"Expone todos los datos",
     ["El PDF muestra RNC, dirección y obligaciones activas.",
      "Sin revelación selectiva para probar solo “está al día”."])
card(s,6.9,4.92,5.55,1.95,GOLD,"Una foto, no un estado en vivo",
     ["Si cae en mora dentro del plazo, el PDF sigue validando.",
      "No interoperable con verificadores internacionales."])
footer(s,num())

# ---- 3. La propuesta ----
s=prs.slides.add_slide(BLANK); header(s,"La propuesta","Una certificación que se verifica sola")
tf=tbox(s,0.9,2.0,11.4,0.7); run(para(tf,first=True,space=0,lh=1.1),"La DGII emite la certificación como una credencial digital firmada que el ciudadano guarda en su billetera.",17,GRAY,italic=True)
card(s,0.9,2.95,5.55,1.75,GREEN,"A prueba de manipulación",["Firmada criptográficamente por la DGII.","Cualquier alteración la invalida al instante."],)
card(s,6.9,2.95,5.55,1.75,NAVY,"Verificable al instante",["Se comprueba sin llamar a la DGII.","Sin papeleo ni intermediarios."])
card(s,0.9,4.95,5.55,1.75,GREEN,"Privacidad por diseño",["El ciudadano revela solo lo necesario.","Controla qué comparte en cada trámite."])
card(s,6.9,4.95,5.55,1.75,NAVY,"Control en tiempo real",["La DGII puede revocarla cuando cambie el estado.","Deja de ser válida en el acto."])
footer(s,num())

# ---- 4. Actores ----
s=prs.slides.add_slide(BLANK); header(s,"El modelo","Un modelo de confianza con tres actores")
card(s,0.9,2.3,3.66,3.5,GREEN,"DGII — Emisor",["Valida el cumplimiento del contribuyente.","Emite y firma la certificación.","Puede revocarla si cambia el estado."],tsize=17,bsize=14)
card(s,4.83,2.3,3.66,3.5,GOLD,"Ciudadano — Titular",["Recibe la credencial en su billetera digital.","La lleva consigo, siempre disponible.","Decide qué datos comparte y con quién."],tsize=17,bsize=14)
card(s,8.76,2.3,3.66,3.5,NAVY,"Banco — Verificador",["Solicita la certificación al cliente.","La valida al instante, sin llamar a la DGII.","Recibe solo lo estrictamente necesario."],tsize=17,bsize=14)
tf=tbox(s,0.9,6.05,11.5,0.5); run(para(tf,first=True,space=0,align=PP_ALIGN.CENTER),"Confianza directa entre las partes, respaldada por la firma de la DGII — sin intermediarios.",14,GRAY,italic=True)
footer(s,num())

# ---- 5. 4 pasos (flujo) ----
s=prs.slides.add_slide(BLANK); header(s,"Cómo funciona","El ciclo de vida, en cuatro pasos")
steps=[("1","Emitir","La DGII valida y emite la certificación firmada.",GREEN),
       ("2","Guardar","El ciudadano la recibe en su billetera digital.",GOLD),
       ("3","Verificar","El banco la valida con revelación selectiva.",NAVY),
       ("4","Revocar","Si cae en mora, la DGII la invalida en el acto.",RED)]
x=0.9; w=2.75; y=2.7; h=2.5; gapx=0.28
for i,(nn,t,d,c) in enumerate(steps):
    cx=x+i*(w+gapx)
    rect(s,cx,y,w,h,LIGHT,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s,cx,y,w,0.7,c,shape=MSO_SHAPE.RECTANGLE)
    tf=tbox(s,cx,y+0.08,w,0.55,anchor=MSO_ANCHOR.MIDDLE); run(para(tf,first=True,space=0,align=PP_ALIGN.CENTER),f"{nn}.  {t}",16,WHITE,bold=True)
    tf=tbox(s,cx+0.3,y+0.95,w-0.55,h-1.1); run(para(tf,first=True,space=0,lh=1.15),d,14,INK)
    if i<3:
        ar=tbox(s,cx+w-0.02,y+h/2-0.25,gapx+0.05,0.5,anchor=MSO_ANCHOR.MIDDLE)
        run(para(ar,first=True,space=0,align=PP_ALIGN.CENTER),"→",22,GRAY,bold=True)
tf=tbox(s,0.9,5.7,11.5,0.5); run(para(tf,first=True,space=0,align=PP_ALIGN.CENTER),"Estándares abiertos de extremo a extremo — emisión, resguardo, verificación y revocación.",14,GRAY,italic=True)
footer(s,num())

# ---- 6. Revelación selectiva ----
s=prs.slides.add_slide(BLANK); header(s,"El diferenciador","Privacidad: revelación selectiva")
tf=tbox(s,0.9,2.0,11.4,0.7); run(para(tf,first=True,space=0,lh=1.1),"Probar que está al día — sin exponer nada más. El ciudadano elige qué revela en cada presentación.",18,NAVY,bold=True)
card(s,0.9,3.0,5.55,2.9,GREEN,"✓  Lo que el banco recibe",["Estado de cumplimiento:  AL DÍA","(y, si se solicita, el RNC)"],tsize=17,bsize=15)
card(s,6.9,3.0,5.55,2.9,GRAY,"✕  Lo que NO se expone",["Razón social","Domicilio fiscal","Obligaciones activas","Régimen y actividad económica"],tsize=17,bsize=15)
footer(s,num())

# ---- 7. Revocación ----
s=prs.slides.add_slide(BLANK); header(s,"El control","Revocación en tiempo real")
bullets(s,0.9,2.15,11.5,[
 ("Si el contribuyente cae en mora, ","la DGII revoca la certificación desde su sistema."),
 ("El rechazo es inmediato. ","Cualquier verificador la da como no válida — sin recuperar ningún documento físico."),
 ("La credencial es una foto firmada. ","No se “edita”: la revocación la invalida, sin volver a exponer al ciudadano."),
 ("En la demostración, ","al revocar, el registro marca al contribuyente EN MORA y el portal del banco muestra “Credencial no válida — estado actual: EN MORA”."),
],size=18,gap=16)
footer(s,num())

# ---- 8. Beneficios ----
s=prs.slides.add_slide(BLANK); header(s,"El valor","Beneficios para todos los actores")
card(s,0.9,2.3,3.66,3.6,GREEN,"Para la DGII",["Menos fraude documental.","Menos carga operativa: menos consultas y validaciones manuales.","Modernización e imagen de identidad digital."],tsize=17,bsize=14)
card(s,4.83,2.3,3.66,3.6,GOLD,"Para el ciudadano",["Privacidad: comparte solo lo necesario.","Control de sus propios datos.","Portátil y disponible 24/7."],tsize=17,bsize=14)
card(s,8.76,2.3,3.66,3.6,NAVY,"Para bancos y entidades",["Confianza instantánea, sin llamadas.","Menor costo y fricción en verificación (KYC).","Menor riesgo de aceptar documentos falsos."],tsize=17,bsize=14)
footer(s,num())

# ---- Demostración del demo (divisor de sección) ----
s=prs.slides.add_slide(BLANK)
rect(s,0,0,SW,7.5,NAVY)
rect(s,0,2.9,SW,0.10,LIME)
logo(s,9.75,0.7,2.7,chip=WHITE)
tf=tbox(s,0.9,1.4,8.4,0.5); run(para(tf,first=True,space=0),"EN VIVO",14,LIME,bold=True)
tf=tbox(s,0.9,1.85,11.4,0.8); run(para(tf,first=True,space=0,lh=1.0),"Demostración del demo",34,WHITE,bold=True,font=FL)
tf=tbox(s,0.9,3.25,11.5,0.9); run(para(tf,first=True,space=0,lh=1.18),"Flujo de punta a punta sobre estándares abiertos: la DGII emite → el ciudadano recibe por correo (QR) → el banco verifica con revelación selectiva → la DGII revoca en tiempo real.",16,WHITE)
for i,(t,d) in enumerate([("1 · Emitir","La DGII valida el RNC y firma la certificación."),
                          ("2 · Tener","El ciudadano la recibe en su wallet (correo + QR)."),
                          ("3 · Verificar","El banco la valida — solo ve “al día”."),
                          ("4 · Revocar","Si cae en mora, deja de ser válida en el acto.")]):
    cx=0.9+i*3.02
    tf=tbox(s,cx,4.55,2.85,1.5)
    run(para(tf,first=True,space=5),t,17,LIME,bold=True)
    run(para(tf,space=0,lh=1.15),d,13,WHITE)
tf=tbox(s,0.9,6.5,11.5,0.5); run(para(tf,first=True,space=0),"Dos portales: Oficina Virtual de la DGII (emisor) y Portal del Banco «Banco de las Antillas» (verificador).",13,LIME,italic=True)

# ---- 9. Qué demostramos ----
s=prs.slides.add_slide(BLANK); header(s,"La prueba de concepto","Lo que ya demostramos, funcionando")
bullets(s,0.9,2.15,11.5,[
 ("Flujo completo de extremo a extremo. ","La DGII emite → el ciudadano recibe (correo y QR) → el banco verifica → la DGII revoca."),
 ("Dos portales. ","Oficina Virtual de la DGII (emisor) y portal del banco (verificador)."),
 ("Privacidad real. ","Revelación selectiva: el banco solo ve “al día”."),
 ("Control real. ","Revocación en tiempo real, con el contraste al día ↔ en mora."),
 ("Estándares abiertos. ","Sin quedar atado a un proveedor; interoperable con billeteras de terceros."),
],size=17,gap=13)
footer(s,num())

# ---- 10. Estándares ----
s=prs.slides.add_slide(BLANK); header(s,"La base tecnológica","Construido sobre estándares abiertos")
bullets(s,0.9,2.15,11.5,[
 ("Formato SD-JWT VC y protocolos OpenID4VCI / OpenID4VP. ","Estándares internacionales para emitir y presentar credenciales."),
 ("Identidad del emisor verificable. ","La DGII se identifica de forma comprobable (did:web) — el verificador confía en el origen."),
 ("Revocación por listas de estado. ","Mecanismo estándar para invalidar credenciales en tiempo real."),
 ("Interoperable y sin dependencia de proveedor. ","Compatible con billeteras y verificadores de terceros; alineado con las tendencias globales de identidad digital."),
],size=17,gap=15)
footer(s,num())

# ---- 11. Hacia producción ----
s=prs.slides.add_slide(BLANK); header(s,"El siguiente paso","Del demo hacia producción")
bullets(s,0.9,2.15,11.5,[
 ("Identidad oficial y registro de confianza. ","Publicar la identidad de la DGII (did:web:dgii.gob.do) y un registro de emisores confiables."),
 ("Firma con HSM y gobierno de claves. ","Nivel de seguridad institucional para la firma de credenciales."),
 ("Integración con los sistemas tributarios. ","Conectar con SIT / OFV como fuente de verdad del cumplimiento."),
 ("Piloto con una entidad financiera. ","Validar adopción, experiencia de usuario e impacto operativo."),
],size=17,gap=15)
footer(s,num())

# ---- 12. Cierre ----
s=prs.slides.add_slide(BLANK)
rect(s,0,0,SW,7.5,NAVY)
rect(s,0,3.05,SW,0.10,LIME)
logo(s,9.75,0.7,2.7,chip=WHITE)                 # logo DGII sobre chip blanco
tf=tbox(s,0.9,1.5,8.4,0.5); run(para(tf,first=True,space=0),"EN RESUMEN",14,LIME,bold=True)
tf=tbox(s,0.9,2.0,11.5,1.1); run(para(tf,first=True,space=0,lh=1.05),"Una certificación fiscal más segura, más privada y verificable al instante.",28,WHITE,bold=True,font=FL)
for i,(t,d) in enumerate([("Menos fraude","Firmada por la DGII, imposible de falsificar."),
                          ("Más privacidad","El ciudadano revela solo lo necesario."),
                          ("Cero fricción","Verificación inmediata y revocación en tiempo real.")]):
    cx=0.9+i*3.95
    tf=tbox(s,cx,3.6,3.7,1.4)
    run(para(tf,first=True,space=4),t,18,LIME,bold=True)
    run(para(tf,space=0,lh=1.15),d,14,WHITE)
tf=tbox(s,0.9,5.7,11.5,0.6); run(para(tf,first=True,space=0),"Gracias.  ¿Preguntas?",22,WHITE,bold=True,font=FL)

out="/home/triageuser/CertVerificableDGII/dgii-demo/Caso-de-uso-DGII-Credencial-Verificable.pptx"
prs.save(out); print("guardado:",out,"·",len(prs.slides._sldIdLst),"slides")
